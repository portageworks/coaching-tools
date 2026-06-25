"""Session Cue Deck — render the worksheet content as a themed .pptx the coach
flips through during a session (and that opens in Google Slides or PowerPoint).

Two stages:
  parse_worksheet(md)  -> a flat list of slide dicts
  build_worksheet_pptx(md, name) -> .pptx bytes

Design follows the app theme (charcoal + slate), safe fonts (Cambria headers,
Calibri body) so it renders true on any machine, and keeps everything visible
on the slide (no speaker notes) — it's a guide the coach reads from, not a
deck presented to the client.
"""
import io
import re
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE

# ── Palette (matches the app: charcoal + slate) ─────────────────────────────────
CHARCOAL   = RGBColor(0x1E, 0x20, 0x22)
CHARCOAL2  = RGBColor(0x3D, 0x41, 0x45)
SLATE      = RGBColor(0x3A, 0x5A, 0x7C)
SLATE_MID  = RGBColor(0x55, 0x79, 0xA0)
SLATE_PALE = RGBColor(0xEE, 0xF4, 0xFB)
INK        = RGBColor(0x2C, 0x30, 0x35)
INK_MID    = RGBColor(0x4A, 0x50, 0x58)
DIM        = RGBColor(0x7A, 0x80, 0x88)
PAPER      = RGBColor(0xFF, 0xFF, 0xFF)
SAY_BG     = RGBColor(0xFB, 0xF6, 0xEC)
SAY_INK    = RGBColor(0x6B, 0x55, 0x1F)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

HEAD_FONT = "Cambria"
BODY_FONT = "Calibri"

SW, SH = Inches(13.333), Inches(7.5)


# ── Parsing ─────────────────────────────────────────────────────────────────────
def _clean(s):
    return re.sub(r"\s+", " ", (s or "").strip())


def parse_worksheet(md):
    """Turn worksheet markdown into an ordered list of slide dicts."""
    # Drop note markers entirely.
    md = re.sub(r"(?m)^[ \t]*\[\[NOTES(?::[a-zA-Z]+)?\]\][ \t]*$", "", md)
    lines = md.split("\n")

    slides = []
    title_name = ""
    section = None            # "before" | "words" | "block"
    block_title = ""
    block_meta = ""
    # current question/story subsection accumulator
    cur = None

    def flush_sub():
        nonlocal cur
        if cur and (cur["saids"] or cur["questions"] or cur["says"] or cur["checks"]):
            if cur["checks"] and not cur["questions"] and not cur["saids"]:
                slides.append({"kind": "checklist", "block": block_title,
                               "title": cur["title"], "items": cur["checks"]})
            else:
                slides.append({"kind": "qa", "block": block_title,
                               "title": cur["title"], "saids": cur["saids"],
                               "questions": cur["questions"], "says": cur["says"]})
        cur = None

    # accumulators for before/words
    anchors = []      # (label, body)
    quotes = []       # (label, body)
    pending_label = None
    pending_body = []
    in_quote = False

    def flush_label():
        nonlocal pending_label, pending_body, in_quote
        if pending_label is not None:
            body = _clean(" ".join(pending_body))
            if section == "before":
                anchors.append((pending_label, body))
            elif section == "words":
                quotes.append((pending_label, body))
        pending_label, pending_body, in_quote = None, [], False

    def flush_section_blocks():
        # emit accumulated before/words slides
        nonlocal anchors, quotes
        if anchors:
            slides.append({"kind": "anchors", "items": anchors})
            anchors = []
        if quotes:
            # chunk quote fields 4 per slide to avoid overflow
            for i in range(0, len(quotes), 4):
                slides.append({"kind": "quotes", "items": quotes[i:i + 4]})
            quotes = []

    i = 0
    while i < len(lines):
        raw = lines[i]
        line = raw.strip()
        i += 1
        if not line:
            continue

        m = re.match(r"^#\s+(.*)", line)
        if m and not line.startswith("##"):
            title_name = re.sub(r"(?i)session worksheet\s*--\s*", "", m.group(1)).strip()
            continue

        m = re.match(r"^##\s+(.*)", line)
        if m and not line.startswith("###"):
            head = m.group(1).strip()
            up = head.upper()
            # close out whatever we were building
            flush_label(); flush_sub()
            if up.startswith("BEFORE YOU START"):
                flush_section_blocks(); section = "before"
            elif up.startswith("IN THEIR WORDS"):
                flush_section_blocks(); section = "words"
            elif up.startswith("SINGLE-DAY") or up.startswith("SINGLE DAY"):
                continue  # subtitle of title; ignore
            elif up.startswith("BLOCK"):
                flush_section_blocks(); section = "block"
                block_title = head
                block_meta = ""
                # divider slide
                slides.append({"kind": "divider", "title": head})
            else:
                flush_section_blocks(); section = "block"
                block_title = head
                slides.append({"kind": "divider", "title": head})
            continue

        m = re.match(r"^###\s+(.*)", line)
        if m:
            flush_sub()
            cur = {"title": m.group(1).strip(), "saids": [], "questions": [],
                   "says": [], "checks": []}
            continue

        # block objective: *...* italic line right under a block
        if section == "block" and cur is None and re.match(r"^\*.+\*$", line):
            block_meta = line.strip("*").strip()
            # attach to last divider
            for s in reversed(slides):
                if s["kind"] == "divider":
                    s["meta"] = block_meta
                    break
            continue

        # bold field label  **Label:**
        m = re.match(r"^\*\*(.+?):\*\*\s*(.*)$", line)
        if m and section in ("before", "words"):
            flush_label()
            pending_label = m.group(1).strip()
            rest = m.group(2).strip()
            if rest:
                pending_body.append(rest)
            continue

        # blockquote line
        if line.startswith(">"):
            q = line.lstrip(">").strip()
            if section in ("before", "words"):
                # skip the "THEY SAID:" marker line itself
                if re.match(r"(?i)^THEY SAID\s*:?$", q):
                    continue
                pending_body.append(q)
            elif cur is not None:
                cur.setdefault("_bq", []).append(q)
            continue

        # checkbox
        m = re.match(r"^-\s*\[ \]\s*(.*)$", line)
        if m and cur is not None:
            cur["checks"].append(m.group(1).strip())
            continue

        # question item
        m = re.match(r"^-\s*(Ask|If surface|If deep):\s*(.*)$", line)
        if m and cur is not None:
            # first, fold any pending blockquote into said cards
            _fold_saids(cur)
            cur["questions"].append((m.group(1), m.group(2).strip()))
            continue

        # say line (italic) inside a subsection
        if cur is not None and re.match(r"^\*.+\*$", line):
            _fold_saids(cur)
            cur["says"].append(re.sub(r"^Say:\s*", "", line.strip("*").strip()))
            continue

        # plain text within before/words label
        if section in ("before", "words") and pending_label is not None:
            pending_body.append(line)
            continue

    # fold trailing
    if cur is not None:
        _fold_saids(cur)
    flush_label(); flush_sub(); flush_section_blocks()

    slides.insert(0, {"kind": "title", "name": title_name or "Client",
                      "date": datetime.now().strftime("%B %d, %Y")})
    return slides


def _fold_saids(cur):
    """Convert accumulated THEY SAID blockquote lines into (label, body) pairs."""
    bq = cur.pop("_bq", None)
    if not bq:
        return
    text = "\n".join(bq)
    segs = re.split(r"THEY SAID\s*(?:--\s*([^:\n]+?))?\s*:", text)
    for j in range(1, len(segs), 2):
        label = _clean(segs[j]) if segs[j] else ""
        body = _clean(segs[j + 1]) if j + 1 < len(segs) else ""
        if body:
            cur["saids"].append((label or "They said", body))


# ── Rendering helpers ───────────────────────────────────────────────────────────
def _txt(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    return tf


def _run(p, text, size, color, font=BODY_FONT, bold=False, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.name = font
    r.font.bold = bold
    r.font.italic = italic
    return r


def _rect(slide, x, y, w, h, fill, line=None, radius=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def _eyebrow(slide, x, y, w, text, color=SLATE_MID):
    tf = _txt(slide, x, y, w, Inches(0.3))
    p = tf.paragraphs[0]
    _run(p, text.upper(), 11, color, font=BODY_FONT, bold=True)
    return tf


# ── Slide builders ──────────────────────────────────────────────────────────────
def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _slide_title(slide, block, title):
    """Top header band (text only — no accent rule)."""
    tf = _txt(slide, Inches(0.6), Inches(0.42), Inches(12.1), Inches(1.0))
    if block:
        p0 = tf.paragraphs[0]
        _run(p0, block.upper(), 12, SLATE_MID, font=BODY_FONT, bold=True)
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    _run(p, title, 30, CHARCOAL, font=HEAD_FONT, bold=True)


def _build_title(s, slide):
    _bg(slide, CHARCOAL)
    tf = _txt(slide, Inches(0.9), Inches(2.6), Inches(11.5), Inches(2.2))
    p = tf.paragraphs[0]
    _run(p, "SESSION CUE", 14, SLATE_MID, font=BODY_FONT, bold=True)
    p2 = tf.add_paragraph(); p2.space_before = Pt(6)
    _run(p2, s["name"], 48, WHITE, font=HEAD_FONT, bold=True)
    p3 = tf.add_paragraph(); p3.space_before = Pt(10)
    _run(p3, "Single-Day Coaching Session  ·  " + s["date"], 16, RGBColor(0x9B, 0xA0, 0xA6))


def _build_divider(s, slide):
    _bg(slide, CHARCOAL)
    tf = _txt(slide, Inches(0.9), Inches(2.7), Inches(11.5), Inches(2.0),
              anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    _run(p, s["title"], 40, WHITE, font=HEAD_FONT, bold=True)
    if s.get("meta"):
        p2 = tf.add_paragraph(); p2.space_before = Pt(12)
        _run(p2, s["meta"], 18, SLATE_MID, italic=True)


def _build_anchors(s, slide):
    _bg(slide, PAPER)
    _slide_title(slide, "", "Before You Start")
    y = Inches(1.7)
    for label, body in s["items"]:
        big = "north star" in label.lower()
        _eyebrow(slide, Inches(0.6), y, Inches(12), label)
        tf = _txt(slide, Inches(0.6), y + Inches(0.34), Inches(12.1), Inches(1.3))
        p = tf.paragraphs[0]
        if big:
            _run(p, body, 26, SLATE, font=HEAD_FONT, bold=True)
            y += Inches(1.5)
        else:
            _run(p, body, 16, INK_MID)
            y += Inches(0.42) + Inches(0.32) * max(1, len(body) // 90 + 1)


def _quote_card(slide, x, y, w, h, label, body):
    _rect(slide, x, y, w, h, SLATE_PALE)
    tf = _txt(slide, x + Inches(0.18), y + Inches(0.12), w - Inches(0.36), h - Inches(0.24))
    p = tf.paragraphs[0]
    _run(p, ("They said · " + label).upper() if label and label != "They said"
         else "They said".upper(), 10, SLATE, font=BODY_FONT, bold=True)
    p2 = tf.add_paragraph(); p2.space_before = Pt(4)
    _run(p2, body, 13, CHARCOAL)


def _build_quotes(s, slide):
    _bg(slide, PAPER)
    _slide_title(slide, "", "In Their Words")
    items = s["items"]
    cols, gap = 2, Inches(0.4)
    cw = (Inches(12.1) - gap) / 2
    ch = Inches(2.3)
    x0, y0 = Inches(0.6), Inches(1.7)
    for idx, (label, body) in enumerate(items):
        r, c = divmod(idx, cols)
        x = x0 + c * (cw + gap)
        y = y0 + r * (ch + Inches(0.25))
        _quote_card(slide, x, y, cw, ch, label, body)


def _build_checklist(s, slide):
    _bg(slide, PAPER)
    _slide_title(slide, s.get("block", ""), s.get("title") or "Checklist")
    tf = _txt(slide, Inches(0.7), Inches(1.9), Inches(11.9), Inches(5.0))
    first = True
    for item in s["items"]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(10)
        _run(p, "☐  ", 20, CHARCOAL2)
        _run(p, item, 18, CHARCOAL)


def _build_qa(s, slide):
    _bg(slide, PAPER)
    _slide_title(slide, s.get("block", ""), s.get("title") or "")
    saids, questions, says = s["saids"], s["questions"], s["says"]
    top = Inches(1.85)
    two = bool(saids) and bool(questions)

    if two:
        lw = Inches(5.3)
        rx = Inches(6.3)
        rw = Inches(6.4)
        _render_saids(slide, Inches(0.6), top, lw, saids)
        _render_questions(slide, rx, top, rw, questions, says)
    elif saids:
        _render_saids(slide, Inches(0.6), top, Inches(12.1), saids)
    else:
        _render_questions(slide, Inches(0.7), top, Inches(11.9), questions, says, big=True)


def _render_saids(slide, x, y, w, saids):
    _eyebrow(slide, x, y, w, "In their words")
    yy = y + Inches(0.38)
    n = len(saids) or 1
    avail = Inches(4.7)
    ch = min(Inches(1.45), Emu(int(avail / n)) - Inches(0.12))
    for label, body in saids:
        _quote_card(slide, x, yy, w, ch, label, body)
        yy = yy + ch + Inches(0.12)


def _render_questions(slide, x, y, w, questions, says, big=False):
    _eyebrow(slide, x, y, w, "Ask")
    tf = _txt(slide, x, y + Inches(0.38), w, Inches(4.6))
    first = True
    qsize = 22 if big else 19
    for typ, text in questions:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(12)
        if typ == "Ask":
            _run(p, "ASK   ", 11, SLATE, font=BODY_FONT, bold=True)
            _run(p, text, qsize, CHARCOAL, bold=True)
        else:
            tag = "SURFACE   " if typ == "If surface" else "DEEP   "
            _run(p, tag, 11, SLATE_MID, font=BODY_FONT, bold=True)
            _run(p, text, qsize - 4, INK_MID)
    for say in says:
        p = tf.add_paragraph(); p.space_before = Pt(8); p.space_after = Pt(4)
        _run(p, "SAY   ", 11, SAY_INK, font=BODY_FONT, bold=True)
        _run(p, say, 15, CHARCOAL2, italic=True)


_BUILDERS = {
    "title": _build_title, "divider": _build_divider, "anchors": _build_anchors,
    "quotes": _build_quotes, "checklist": _build_checklist, "qa": _build_qa,
}


def build_worksheet_pptx(worksheet_md, client_name):
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    blank = prs.slide_layouts[6]
    slides = parse_worksheet(worksheet_md)
    # ensure the title uses the real client name if the worksheet omitted it
    if slides and slides[0]["kind"] == "title" and client_name:
        slides[0]["name"] = client_name
    for s in slides:
        slide = prs.slides.add_slide(blank)
        _BUILDERS.get(s["kind"], lambda *_: None)(s, slide)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
